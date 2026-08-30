"""Saathi pitch deck. Every slide carries an image. Swiss grotesque, huge type,
few words, alternating dark full-bleed and light split panels.

Run:  /tmp/deckenv/bin/python build_deck.py
"""

from pathlib import Path

from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ASSETS = Path(
    "/Users/abhinav.singh-ext/.cursor/projects/"
    "Users-abhinav-singh-ext-Documents-codes/assets"
)
TMP = Path("/tmp/deck_crops")
TMP.mkdir(exist_ok=True)
OUT = Path(__file__).parent / "Saathi.pptx"
SW, SH = Inches(13.333), Inches(7.5)

INK = RGBColor(0x0B, 0x0B, 0x0C)
PAPER = RGBColor(0xF4, 0xF2, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BONE = RGBColor(0xD8, 0xD5, 0xCE)
MUTE_D = RGBColor(0x8E, 0x8D, 0x91)
MUTE_L = RGBColor(0x6E, 0x6C, 0x69)
FAINT = RGBColor(0xBA, 0xB7, 0xB1)
ORANGE = RGBColor(0xF0, 0x3C, 0x0A)

F = "Arial"

# real hardware, shot at the hackathon
RIG_ROPE = "WhatsApp_Image_2026-08-29_at_14.29.23-dfccab0b-7289-4270-8bb9-bfe2282d6ab4.jpg"
RIG_VIEWER = "WhatsApp_Image_2026-08-29_at_14.29.21-0489d5da-83f7-49f6-8f62-ff65c5f8bcd5.jpg"

prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
M = Inches(1.0)


def slide(bg=INK):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid()
    r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def alpha(shape, opacity):
    solid = shape._element.spPr.find(qn("a:solidFill"))
    clr = solid.find(qn("a:srgbClr"))
    el = etree.SubElement(clr, qn("a:alpha"))
    el.set("val", str(int(opacity * 1000)))


def cover(s, name, left=0, top=0, width=None, height=None):
    width, height = width or SW, height or SH
    p = ASSETS / name
    iw, ih = Image.open(p).size
    if iw / ih > width / height:
        h, w = height, Emu(int(height * iw / ih))
    else:
        w, h = width, Emu(int(width * ih / iw))
    return s.shapes.add_picture(
        str(p), Emu(int(left + (width - w) / 2)), Emu(int(top + (height - h) / 2)), w, h
    )


def boxed(s, name, left, top, width, height):
    """Hard-crop the source to the box aspect so nothing overflows."""
    src = ASSETS / name
    dst = TMP / f"box_{Path(name).stem}_{int(width)}_{int(height)}.png"
    if not dst.exists():
        im = Image.open(src).convert("RGB")
        iw, ih = im.size
        target = width / height
        if iw / ih > target:
            nw = int(ih * target)
            im = im.crop(((iw - nw) // 2, 0, (iw - nw) // 2 + nw, ih))
        else:
            nh = int(iw / target)
            im = im.crop((0, (ih - nh) // 2, iw, (ih - nh) // 2 + nh))
        im.save(dst)
    return s.shapes.add_picture(str(dst), left, top, width, height)


def square(s, name, left, top, size_in):
    """Centre-crop to a square and place it."""
    src = ASSETS / name
    dst = TMP / f"sq_{Path(name).stem}.png"
    if not dst.exists():
        im = Image.open(src).convert("RGB")
        w, h = im.size
        k = min(w, h)
        im.crop(((w - k) // 2, (h - k) // 2, (w - k) // 2 + k, (h - k) // 2 + k)).resize(
            (700, 700), Image.LANCZOS
        ).save(dst)
    return s.shapes.add_picture(str(dst), left, top, Inches(size_in), Inches(size_in))


def scrim(s, opacity, left=0, top=0, width=None, height=None, color=INK):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width or SW, height or SH)
    r.fill.solid()
    r.fill.fore_color.rgb = color
    r.line.fill.background()
    r.shadow.inherit = False
    alpha(r, opacity)
    return r


def panel(s, left, width, color=PAPER):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, 0, width, SH)
    r.fill.solid()
    r.fill.fore_color.rgb = color
    r.line.fill.background()
    r.shadow.inherit = False
    return r


def t(s, body, left, top, width, size=16, color=BONE, bold=False, font=F,
      space=1.08, caps=False, spacing=0, align=PP_ALIGN.LEFT):
    box = s.shapes.add_textbox(left, top, width, Inches(0.6))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(body.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = space
        run = p.add_run()
        run.text = line.upper() if caps else line
        f = run.font
        f.name, f.size, f.bold = font, Pt(size), bold
        f.color.rgb = color
        run.font._rPr.set("spc", str(int(spacing * 100)))
    return box


def big(s, body, left, top, width, size=72, color=WHITE, space=1.0):
    return t(s, body, left, top, width, size=size, color=color, bold=True,
             space=space, spacing=-1.4)


def kicker(s, label, left=M, top=Inches(0.78), color=ORANGE):
    t(s, label, left, top, Emu(int(SW - left - Inches(0.6))), size=10.5,
      color=color, bold=True, caps=True, spacing=2.6)


def bar(s, left, top, width=Inches(2.6), color=ORANGE, thick=Pt(4)):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, thick)
    r.fill.solid()
    r.fill.fore_color.rgb = color
    r.line.fill.background()
    r.shadow.inherit = False
    return r


# ══════════════════════════ 1 · Title
s = slide()
cover(s, "saathi-hero-hands.png")
scrim(s, 40)
big(s, "Saathi", M, Inches(2.45), Inches(9), size=104)
t(s, "A robot whose only job is to hold a hand in the dark.",
  M, Inches(4.36), Inches(10), size=27, color=WHITE, space=1.15, spacing=-0.4)
bar(s, M, Inches(5.32), Inches(1.5))
t(s, "साथी  ·  companion, in Nepali", M, Inches(5.62), Inches(8), size=14, color=FAINT, spacing=1.2)
t(s, "Himalaya Robotics Hack 2026   ·   Track 2, Action   ·   Voice",
  M, Inches(6.72), Inches(9), size=10, color=MUTE_D, caps=True, spacing=2.2)

# ══════════════════════════ 1b · Why me  (the only colour photo in the deck)
s = slide(PAPER)
cover(s, "saathi-personal.jpg", left=Inches(6.9), width=Inches(6.433))
panel(s, 0, Inches(6.9))
kicker(s, "Why this one")
big(s, "Part of my family\nis from these\nmountains.",
    M, Inches(1.95), Inches(5.6), size=46, color=INK, space=1.08)
bar(s, M, Inches(4.55), Inches(4.9), color=INK, thick=Pt(1))
t(s, "I have ridden these roads. When the Bhote Koshi came\n"
     "down it went through places I know.",
  M, Inches(4.82), Inches(5.4), size=15, color=MUTE_L, space=1.5)
big(s, "This is not a case study\nI picked.",
    M, Inches(5.98), Inches(5.6), size=25, color=INK, space=1.14)

# ══════════════════════════ 2 · The event
s = slide()
cover(s, "saathi-valley-destroyed.png")
scrim(s, 26)
kicker(s, "Rasuwa, Nepal · 26 August 2026")
big(s, "A glacier fell\n1,200 metres.", M, Inches(3.9), Inches(9), size=76, space=1.02)
t(s, "It dammed a river. The dam burst.", M, Inches(6.1), Inches(9), size=22, color=WHITE, spacing=-0.3)

# ══════════════════════════ 3 · Numbers  (photo + light panel)
s = slide(PAPER)
cover(s, "saathi-missing-wall.png", width=Inches(6.5))
panel(s, Inches(6.5), Inches(6.833))
L = Inches(7.35)
kicker(s, "Four days later", left=L)
pairs = [("669", "confirmed dead"), ("2,900", "still missing"),
         ("93,000", "affected"), ("42 km", "of road, every bridge")]
y = Inches(1.62)
for n, lab in pairs:
    big(s, n, L, y, Inches(3.0), size=42, color=INK)
    t(s, lab, Emu(int(L + Inches(2.55))), Emu(int(y + Inches(0.24))), Inches(3.0),
      size=12, color=MUTE_L, caps=True, spacing=1.8)
    y = Emu(int(y + Inches(0.94)))
bar(s, L, Inches(5.5), Inches(4.7), color=INK, thick=Pt(1))
big(s, "Nepal asked for help\nwith two things.",
    L, Inches(5.78), Inches(5.0), size=25, color=INK, space=1.14)
t(s, "Tunnel rescue. Identifying the dead.", L, Inches(6.72), Inches(5.0),
  size=15, color=ORANGE, bold=True)

# ══════════════════════════ 4 · The gap
s = slide()
cover(s, "saathi-human-cost.png")
scrim(s, 46)
kicker(s, "The gap nothing occupies")
big(s, "The warning is minutes.", M, Inches(3.5), Inches(11.4), size=54)
big(s, "The rescue is hours.", M, Inches(4.5), Inches(11.4), size=54, color=ORANGE)

# ══════════════════════════ 5 · Die alone  (light panel + photo)
s = slide(PAPER)
cover(s, "saathi-empty-void.png", left=Inches(6.9), width=Inches(6.433))
panel(s, 0, Inches(6.9))
kicker(s, "What actually kills them")
big(s, "They do not die\nof injury.", M, Inches(2.15), Inches(5.6), size=48, color=INK, space=1.08)
big(s, "They die alone.", M, Inches(3.85), Inches(5.6), size=48, color=ORANGE)
bar(s, M, Inches(4.95), Inches(4.9), color=INK, thick=Pt(1))
t(s, "So doctrine is to get a voice to the person and hold\n"
     "it there for the whole dig. Today that means a human\n"
     "face down in silt, one arm through a gap, six hours.",
  M, Inches(5.22), Inches(5.4), size=15, color=MUTE_L, space=1.5)

# ══════════════════════════ 6 · Reveal
s = slide()
cover(s, "saathi-product.png")
scrim(s, 20)
scrim(s, 74, width=Inches(5.7))
kicker(s, "Introducing")
big(s, "Saathi", M, Inches(2.4), Inches(4.8), size=80)
bar(s, M, Inches(3.9), Inches(1.5))
t(s, "It enters the gap a body\ncannot. It finds a hand.\nIt holds on.",
  M, Inches(4.25), Inches(4.6), size=24, color=WHITE, space=1.3, spacing=-0.4)
t(s, "Company, not a crane.", M, Inches(6.15), Inches(4.6), size=15, color=MUTE_D)

# ══════════════════════════ 7 · Four jobs  (four image thumbnails)
s = slide(PAPER)
kicker(s, "What the hand is for")
jobs = [
    ("saathi-hero-hands.png", "HOLD", "Closes on force,\nnot on position."),
    ("saathi-water.png", "SUSTAIN", "Water, light, air\ndown the same tether."),
    ("saathi-tether.png", "SPEAK", "Nepali, generated\non the device."),
    ("saathi-marker.png", "MARK", "Records the place.\nMoves nothing."),
]
x = M
for img, title, sub in jobs:
    square(s, img, x, Inches(1.55), 2.5)
    big(s, title, x, Inches(4.28), Inches(2.6), size=27, color=INK)
    t(s, sub, x, Inches(4.86), Inches(2.5), size=13, color=MUTE_L, space=1.45)
    x = Emu(int(x + Inches(2.9)))
bar(s, M, Inches(6.16), Inches(11.3), color=INK, thick=Pt(1))
big(s, "One hand. Four jobs. All of them contact.",
    M, Inches(6.44), Inches(11.4), size=32, color=INK)

# ══════════════════════════ 8 · Force
s = slide()
cover(s, "saathi-hero-hands.png")
scrim(s, 54)
kicker(s, "The engineering")
big(s, "It closes on force.\nNot on position.", M, Inches(2.5), Inches(11), size=62, space=1.06)
t(s, "A normal gripper is told where to stop. Give it a hand instead of a box and it closes anyway.\n"
     "Saathi is told how hard it may squeeze, and stops wherever that happens to be.",
  M, Inches(4.85), Inches(11.2), size=16, color=BONE, space=1.5)
t(s, "ROBSTRIDE ACTUATORS  ·  MIT PROTOCOL  ·  TORQUE AT 30 HZ  ·  THE DRIVER DECODES IT,\n"
     "THE ROBOT CLASS DISCARDED IT, WE PUT IT BACK",
  M, Inches(6.2), Inches(11.2), size=11, color=ORANGE, space=1.6, spacing=1.4, bold=True)

# ══════════════════════════ 9 · Offline  (dark night photo)
s = slide()
cover(s, "saathi-no-signal.png")
scrim(s, 42)
kicker(s, "It has to work with nothing")
big(s, "No internet\nin that valley.", M, Inches(1.75), Inches(9), size=64, space=1.04)
t(s, "Every cloud assistant on earth is inert down there. So the speech runs on the robot.",
  M, Inches(4.16), Inches(11), size=17, color=BONE)
models = [("Kriti", "Nepali ASR  ·  rank 1  ·  MIT  ·  built in Kathmandu"),
          ("Kala", "Nepali TTS  ·  60 MB ONNX  ·  50x real time on CPU"),
          ("Voices v0", "419 speakers  ·  match a family member")]
y = Inches(4.9)
for name, spec in models:
    big(s, name, M, y, Inches(3.2), size=20)
    t(s, spec, Inches(4.0), Emu(int(y + Inches(0.07))), Inches(8), size=12.5,
      color=FAINT, spacing=0.6)
    y = Emu(int(y + Inches(0.6)))
t(s, "About 2 GB. On battery. Underground.", M, Inches(6.72), Inches(11), size=15,
  color=ORANGE, bold=True)

# ══════════════════════════ 10 · Tether
s = slide()
cover(s, "saathi-tether.png")
scrim(s, 34)
kicker(s, "Why it is on a wire")
big(s, "Radio does not pass\nthrough rubble.", M, Inches(3.35), Inches(11), size=54, space=1.06)
t(s, "So the cable is the product. Power in, so it never dies mid-hold. Voices in and out,\n"
     "so a doctor in Kathmandu can reach the bottom of the hole.",
  M, Inches(5.5), Inches(11.2), size=16, color=BONE, space=1.5)

# ══════════════════════════ 11 · Sustain
s = slide()
cover(s, "saathi-water.png")
scrim(s, 36)
kicker(s, "02 · Sustain")
big(s, "The hold is why\nthey keep answering.", M, Inches(3.1), Inches(11), size=50, space=1.08)
big(s, "The water is why they are still alive to answer.",
    M, Inches(4.95), Inches(11.4), size=26, color=ORANGE)
t(s, "Same gripper. Same tether. No second mechanism.",
  M, Inches(6.15), Inches(11), size=15, color=BONE)

# ══════════════════════════ 12 · Mark
s = slide()
cover(s, "saathi-marker.png")
scrim(s, 40)
kicker(s, "04 · Mark")
big(s, "Sometimes it\narrives too late.", M, Inches(2.6), Inches(11), size=56, space=1.06)
t(s, "Rasuwa is Tamang and Buddhist. Downstream is Hindu. The rites differ at each end of the\n"
     "same river, and neither is ours to perform. So it records the position, moves nothing, and\n"
     "leaves a marker so the family and the right people can find the place.",
  M, Inches(4.85), Inches(11.4), size=16, color=BONE, space=1.5)

# ══════════════════════════ 13 · Built this weekend  (real hardware)
s = slide(PAPER)
kicker(s, "Honest scope · shot at this table")
big(s, "Built this weekend.", M, Inches(1.28), Inches(11), size=44, color=INK)
boxed(s, RIG_ROPE, M, Inches(2.22), Inches(5.42), Inches(2.55))
boxed(s, RIG_VIEWER, Inches(6.9), Inches(2.22), Inches(5.42), Inches(2.55))
t(s, "The arm, the rope rig, the e-stop", M, Inches(4.9), Inches(5.4), size=10,
  color=MUTE_L, caps=True, spacing=1.6)
t(s, "Recorded episodes in the LeRobot viewer", Inches(6.9), Inches(4.9), Inches(5.4),
  size=10, color=MUTE_L, caps=True, spacing=1.6)
bar(s, M, Inches(5.3), Inches(11.32), color=INK, thick=Pt(1))
lines = [
    ("REAL", "6-DOF RobStride arm live on LeRobot, MIT protocol over CAN."),
    ("REAL", "A force-limited close, replayed on a live human hand, Nepali speech on contact."),
    ("REAL", "Per-joint torque pulled out of a driver that was discarding it."),
    ("RENDER", "The tracked chassis. Everything above was shot on this table."),
]
y = Inches(5.56)
for tag, val in lines:
    c = ORANGE if tag == "REAL" else MUTE_L
    t(s, tag, M, y, Inches(1.1), size=9.5, color=c, bold=True, spacing=2)
    t(s, val, Inches(2.15), Emu(int(y - Inches(0.04))), Inches(10.3), size=13, color=INK)
    y = Emu(int(y + Inches(0.46)))

# ══════════════════════════ 14 · Close
s = slide()
cover(s, "saathi-in-void.png")
scrim(s, 44)
big(s, "It stays.", M, Inches(3.05), Inches(8), size=104)
bar(s, M, Inches(4.95), Inches(1.5))
t(s, "साथी  ·  Saathi  ·  no one waits alone",
  M, Inches(5.3), Inches(9), size=16, color=FAINT, spacing=1.4)

prs.save(OUT)
print(f"wrote {OUT} · {len(prs.slides._sldIdLst)} slides")
